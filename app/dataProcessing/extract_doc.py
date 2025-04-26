import os
from pathlib import Path
import re
import fitz
import docx
import pptx
import spacy
import pymongo
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from typing import List, Dict

# === PARAMÈTRES ===
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
os.environ["TESSDATA_PREFIX"] = r"C:\Program Files\Tesseract-OCR\tessdata"
nlp = spacy.load("fr_core_news_md")

# === Notions-clés ===
NOTIONS_CLES = {
    "dérivée": ["dérivée", "fonction", "tangente", "variation"],
    "énergie": ["énergie", "masse", "relativité", "vitesse"],
    "gaz parfait": ["gaz", "pression", "volume", "température"],
    "cinématique": ["vitesse", "accélération", "mouvement"],
    "chimie": ["molécule", "atome", "liaison", "réaction"],
    "électricité": ["tension", "courant", "résistance"]
}

# === OCR Image ===
def ocr_image(image: Image.Image) -> str:
    return pytesseract.image_to_string(image, lang='fra')

# === OCR PDF ===
def ocr_pdf_scans(pdf_path: str) -> str:
    images = convert_from_path(pdf_path)
    texte = ""
    for img in images:
        texte += ocr_image(img) + "\n"
    return texte

# === Extraction d'équations simples via OCR + Regex ===
def extraire_equations_simples(texte: str) -> List[str]:
    motif = r"[A-Za-z0-9\(\)\^\*\+\-/=\[\]{},.\s]{3,}"
    candidats = re.findall(motif, texte)
    equations = [c.strip() for c in candidats if any(op in c for op in ['=', '+', '-', '*', '/']) and len(c) <= 100]
    return list(set(equations))

# === Extraction PDF ===
def extraire_pdf(fichier):
    doc = fitz.open(fichier)
    texte, images = "", []
    for page in doc:
        texte += page.get_text("text") + "\n"
        for img in page.get_images(full=True):
            base_image = doc.extract_image(img[0])
            images.append(base_image["image"])
    if not texte.strip():
        texte = ocr_pdf_scans(fichier)
    return texte.strip(), images

# === Extraction DOCX ===
def extraire_docx(fichier):
    doc = docx.Document(fichier)
    texte = "\n".join(p.text for p in doc.paragraphs)
    return texte.strip(), []

# === Extraction PPTX ===
def extraire_pptx(fichier):
    prs = pptx.Presentation(fichier)
    texte = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texte += shape.text + "\n"
    return texte.strip(), []

# === Traitement NLP ===
def detecter_notion(contexte):
    doc = nlp(contexte)
    tokens = [token.lemma_.lower() for token in doc if token.pos_ == "NOUN"]
    scores = {notion: sum(1 for mot in mots if mot in tokens) for notion, mots in NOTIONS_CLES.items()}
    return max(scores, key=scores.get) if max(scores.values()) > 0 else "notion inconnue"

def enrichir_formules(formules, texte):
    enrichies = []
    for f in formules:
        if f not in texte: continue
        contexte = texte[max(0, texte.find(f)-100):texte.find(f)+100]
        notion = detecter_notion(contexte)
        enrichies.append({"formule": f, "contexte": contexte, "notion_associee": notion})
    return enrichies

# === Traitement dossier récursif avec détection niveau/matière ===
def traiter_dossier(path_dossier):
    documents = []
    for root, _, fichiers in os.walk(path_dossier):
        for fichier in fichiers:
            if not fichier.endswith((".pdf", ".docx", ".pptx")):
                continue
            chemin = os.path.join(root, fichier)
            rel_path = Path(chemin).relative_to(path_dossier)
            niveau = rel_path.parts[0] if len(rel_path.parts) > 1 else "inconnu"
            matiere = rel_path.parts[1] if len(rel_path.parts) > 2 else "inconnu"

            print(f"🔍 Traitement : {chemin}")
            ext = Path(fichier).suffix.lower()
            if ext == ".pdf": texte, images = extraire_pdf(chemin)
            elif ext == ".docx": texte, images = extraire_docx(chemin)
            elif ext == ".pptx": texte, images = extraire_pptx(chemin)
            else: continue

            equations_detectees = extraire_equations_simples(texte)
            formules_enrichies = enrichir_formules(equations_detectees, texte)

            documents.append({
                "nom": fichier,
                "texte": texte,
                "formules": formules_enrichies,
                "equations_latex": [],
                "nb_formules": len(formules_enrichies),
                "nb_images": len(images),
                "niveau": niveau,
                "matiere": matiere
            })
    return documents

# === Enregistrement MongoDB ===
def sauvegarder_mongodb(documents):
    client = pymongo.MongoClient("mongodb://localhost:27017")
    db = client["cours_lycee"]
    col = db["documents"]
    col.insert_many(documents)
    print("✅ Documents enregistrés dans MongoDB")

# === MAIN ===
if __name__ == "__main__":
    dossier = r"C:\Users\Romaissae\OneDrive\Desktop\data"
    docs = traiter_dossier(dossier)
    sauvegarder_mongodb(docs)
    print("📆 Pipeline terminé avec", len(docs), "documents.")
